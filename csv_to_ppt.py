import os
import shutil
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from datetime import datetime
from tkinter import Tk
from tkinter.filedialog import askopenfilename

# Ouvrir une boîte de dialogue pour sélectionner le fichier CSV
Tk().withdraw()  # Cacher la fenêtre Tkinter principale
file_path = askopenfilename(title="Sélectionner le fichier CSV", filetypes=[("CSV files", "*.csv")])

# Charger les données du fichier CSV
df = pd.read_csv(file_path, sep=';')

# Obtenir la date actuelle et formater le nom du répertoire
current_date = datetime.now().strftime("%Y-%m-%d")
directory_name = f"{current_date} - Thursday's Meeting"

# Créer le répertoire s'il n'existe pas déjà
if not os.path.exists(directory_name):
    os.makedirs(directory_name)

# Copier le fichier CSV sélectionné vers le nouveau répertoire
csv_dest_path = os.path.join(directory_name, os.path.basename(file_path))
shutil.copy(file_path, csv_dest_path)

# Charger la présentation à partir du modèle
template_path = r'C:\Users\natha\Desktop\vs\TEMPLATE.pptx'  # Mettre à jour le chemin si nécessaire
prs = Presentation(template_path)

########################################## Fonction pour supprimer les zones de texte vides
def remove_empty_textboxes(slide):
    for shape in slide.shapes:
        if shape.has_text_frame and not shape.text.strip():
            sp = shape.element
            sp.getparent().remove(sp)

# Fonction pour ajouter la date sur une diapositive
def add_date_to_slide(slide):
    date_str = datetime.now().strftime("%Y-%m-%d")  # Format de la date
    textbox = slide.shapes.add_textbox(Cm(17.78), Cm(1.27), Cm(5.08), Cm(1.27))  # Utiliser des cm ici
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = date_str
    p.font.size = Pt(12)  # Taille de la police de la date

# Fonction pour ajouter le numéro de page sur une diapositive
def add_slide_number(slide, number):
    textbox = slide.shapes.add_textbox(Cm(33), Cm(17.27), Cm(2.54), Cm(1.27))  # Utiliser des cm pour la position en bas à droite
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = str(number)
    p.font.size = Pt(12)  # Taille de la police du numéro de page
    p.font.color.rgb = RGBColor(0, 0, 0)  # Couleur du texte

def delete_slide(prs, slide_index):
    xml_slides = prs.slides._sldIdLst  # Obtenir la liste XML des diapositives
    slides = list(xml_slides)  # Convertir en liste Python
    
    # Vérifier si l'index est dans les limites
    if 0 <= slide_index < len(slides):
        xml_slides.remove(slides[slide_index])  # Supprimer la diapositive à l'index donné
    else:
        print(f"Index {slide_index} hors limites.")

def delete_slide(prs, slide_index):
    xml_slides = prs.slides._sldIdLst  # Obtenir la liste XML des diapositives
    slides = list(xml_slides)  # Convertir en liste Python
    
    # Vérifier si l'index est dans les limites
    if 0 <= slide_index < len(slides):
        xml_slides.remove(slides[slide_index])  # Supprimer la diapositive à l'index donné
    else:
        print(f"Index {slide_index} hors limites.")


# Ajouter une diapositive pour chaque équipe
# Fonction pour ajouter une diapositive pour chaque équipe
# Fonction pour ajouter une diapositive pour chaque équipe
def add_team_slide(prs, team_name, team_data):
    # Déterminer le titre en fonction du nom de l'équipe
    if team_name == "Mathieu Palu":
        title_text = "PB&D's Trades"
    elif team_name == "Claire Bernard":
        title_text = "Bank's Trades"
    else:
        title_text = f"Trades de l'équipe : {team_name}"
    
    # Ajouter une diapositive avec un titre (utiliser le layout du modèle)
    slide_layout = prs.slide_layouts[5]  # Diapositive vide
    slide = prs.slides.add_slide(slide_layout)
    
    # Ajouter le titre de la diapositive
    title = slide.shapes.title
    title.text = title_text
    
    # Définir la position du tableau sur la diapositive
    x, y, cx, cy = Cm(1.27), Cm(5), Cm(21.59), Cm(12.7)  # Utilisation de cm ici
    
    # Déterminer le nombre de lignes et de colonnes pour le tableau
    rows, cols = len(team_data) + 1, len(team_data.columns)
    
    # Ajouter un tableau
    table = slide.shapes.add_table(rows, cols, x, y, cx, cy).table
    
    # Ajouter les bordures du tableau
    set_table_borders(table)
    
    # Continuer à définir les données du tableau (comme dans ton code initial)
    # Ajouter les en-têtes de colonnes
    font_size = Pt(10)
    green_color = RGBColor(144, 238, 144)  # Vert clair
    white_color = RGBColor(255, 255, 255)  # Blanc
    black_color = RGBColor(0, 0, 0)  # Noir

    for col_idx, col_name in enumerate(team_data.columns):
        cell = table.cell(0, col_idx)
        cell.text = col_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = green_color  # Fond vert clair
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = font_size
                run.font.color.rgb = white_color  # Texte blanc

    # Ajouter les données du tableau
    for row_idx, trade in enumerate(team_data.values):
        for col_idx, value in enumerate(trade):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # Fond blanc pour les lignes de données
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = font_size
                    run.font.color.rgb = black_color  # Texte noir

    return slide
# Fonction utilitaire pour définir les bordures des cellules
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

# Fonction utilitaire pour définir les bordures des cellules
# Fonction utilitaire pour définir les bordures des cellules
def set_table_borders(table):
    for row in table.rows:
        for cell in row.cells:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            
            # Ajouter une bordure noire
            border_xml = (
                '<a:lnL w="12700" cap="flat" cmpd="sng" algn="ctr" {0}><a:solidFill><a:srgbClr val="000000"/></a:solidFill></a:lnL>'
                '<a:lnR w="12700" cap="flat" cmpd="sng" algn="ctr" {0}><a:solidFill><a:srgbClr val="000000"/></a:solidFill></a:lnR>'
                '<a:lnT w="12700" cap="flat" cmpd="sng" algn="ctr" {0}><a:solidFill><a:srgbClr val="000000"/></a:solidFill></a:lnT>'
                '<a:lnB w="12700" cap="flat" cmpd="sng" algn="ctr" {0}><a:solidFill><a:srgbClr val="000000"/></a:solidFill></a:lnB>'
            ).format(nsdecls('a'))

            borders = parse_xml(border_xml)
            tcPr.append(borders)
# Liste des équipes uniques
teams = df['Equipe'].unique()
print(teams)

# Ajouter une diapositive pour chaque équipe
slides_to_add = []

for team in teams:
    team_data = df[df['Equipe'] == team]
    slide = add_team_slide(prs, team, team_data)
    slides_to_add.append(slide)

# Réorganiser les slides pour les insérer entre la première et la deuxième
for slide in slides_to_add:
    prs.slides._sldIdLst.insert(1, prs.slides._sldIdLst[-1])

# Ajouter la date à la première diapositive
add_date_to_slide(prs.slides[0])

# Ajouter le numéro de page à chaque diapositive
for i, slide in enumerate(prs.slides):
    add_slide_number(slide, i + 1)

################################ Supprimer les zones de texte vides pour toutes les diapositives
for slide in prs.slides:
    remove_empty_textboxes(slide)
# Suppression de la 2ème diapositive (index 1 car l'index est basé sur 0)
delete_slide(prs, 1)
# Sauvegarder la présentation PowerPoint dans le même répertoire que le fichier CSV
output_path = os.path.join(directory_name, 'pipeline_presentation.pptx')
prs.save(output_path)

print(f"Présentation PowerPoint et fichier CSV enregistrés dans le répertoire : {directory_name}")
